"""Analyze PPTX template structure — find all shapes, text, positions."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

prs = Presentation(r'c:\Users\Lenovo\OneDrive - TMI\Documents\TMI related\AWS Summit\2026\Automatic Reporting\bedrock-agentcore-solution\ESG Document\ESG Performance Dashboard update 3.pptx')

print(f'Slide count: {len(prs.slides)}')
print(f'Slide width: {prs.slide_width} ({prs.slide_width / 914400:.2f} inches)')
print(f'Slide height: {prs.slide_height} ({prs.slide_height / 914400:.2f} inches)')
print()

for i, slide in enumerate(prs.slides):
    print(f'=== SLIDE {i+1} ===')
    for j, shape in enumerate(slide.shapes):
        txt = ''
        if shape.has_text_frame:
            txt = shape.text_frame.text[:100].replace('\n', ' | ')
        shape_type = str(shape.shape_type).split('(')[0] if '(' in str(shape.shape_type) else str(shape.shape_type)
        print(f'  [{j:2d}] type={shape_type:<20} name="{shape.name}"')
        print(f'       pos=({shape.left/914400:.2f}", {shape.top/914400:.2f}") size=({shape.width/914400:.2f}" x {shape.height/914400:.2f}")')
        if txt:
            print(f'       text="{txt}"')
    print()
