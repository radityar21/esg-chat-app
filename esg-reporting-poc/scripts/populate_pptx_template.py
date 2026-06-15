"""
Populate ESG PPTX Template — v2 (fixed charts, matching template style)
Run: python populate_pptx_template.py
Output: ESG_Executive_Summary_FILLED.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from io import BytesIO
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np

TEMPLATE_PATH = r'c:\Users\Lenovo\OneDrive - TMI\Documents\TMI related\AWS Summit\2026\Automatic Reporting\bedrock-agentcore-solution\ESG Document\ESG Performance Dashboard update 3.pptx'
OUTPUT_PATH = r'c:\Users\Lenovo\OneDrive - TMI\Documents\TMI related\AWS Summit\2026\Automatic Reporting\bedrock-agentcore-solution\ESG Document\ESG_Executive_Summary_FILLED.pptx'

# Text replacements per slide (shape_name -> new_text)
SLIDE1_TEXT = {
    "Text 5": "22.03M", "Text 6": "tCO\u2082e  \u25bc 4.07%",
    "Text 12": "239.4", "Text 13": "tCO\u2082e/IDR Bn  \u25bc 11.3%",
    "Text 19": "3.36", "Text 20": "/5.0",
    "Text 26": "99.78%", "Text 27": "of total footprint",
    "Text 41": "Stable", "Text 42": "Emissions -4.07% YoY",
    "Text 48": "Strong", "Text 49": "Female 42.3%, Training 49.3h",
    "Text 55": "Gap", "Text 56": "IFRS/ESRS not established",
    "Text 64": "Establish sector emission targets (O&G, Cement, Steel)",
    "Text 67": "Upgrade PCAF data quality 3.36 \u2192 \u22643.0",
    "Text 70": "Obtain limited assurance (Scope 1, 2, 3)",
}
SLIDE2_TEXT = {
    "Text 16": "99.78% of total emissions from financed activities (Scope 3 Cat 15). Top 3 sectors = 85.93% of financed emissions despite only 20.70% of portfolio.",
    "Text 20": "5", "Text 25": "3.36/5", "Text 30": "-4.07%",
}
SLIDE3_TEXT = {
    "Text 9": "24,997", "Text 14": "42.32%", "Text 19": "49.3 hrs",
    "Text 32": "Partial (6 gaps)", "Text 41": "Partial (6 gaps)", "Text 50": "Partial (3 gaps)",
}


# === CHART FUNCTIONS (fixed sizing, no overlap, matching template palette) ===

def chart_stacked_bar():
    fig, ax = plt.subplots(figsize=(4.2, 2.4))
    fig.patch.set_facecolor('white')
    s1, s2, s3 = [3324, 3403], [44601, 44614], [22930000, 21976797]
    x = np.arange(2)
    ax.bar(x, s3, 0.45, label='Scope 3 Cat 15 (Financed)', color='#1B5E20')
    ax.bar(x, s2, 0.45, bottom=s3, label='Scope 2 (Market-based)', color='#43A047')
    ax.bar(x, s1, 0.45, bottom=[a+b for a,b in zip(s3,s2)], label='Scope 1 (Direct)', color='#A5D6A7')
    for i, t in enumerate([a+b+c for a,b,c in zip(s1,s2,s3)]):
        ax.annotate(f'{t/1e6:.2f}M', (i,t), xytext=(0,4), textcoords="offset points", ha='center', fontsize=8, fontweight='bold', color='#1B3A6B')
    ax.set_xticks(x); ax.set_xticklabels(['FY2023','FY2024'], fontsize=9)
    ax.set_title('Total GHG Emissions', fontsize=10, fontweight='bold', color='#1B3A6B', pad=6)
    ax.legend(fontsize=6, loc='upper center', bbox_to_anchor=(0.5,-0.12), ncol=3, frameon=False)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,p: f'{v/1e6:.0f}M'))
    ax.grid(axis='y', alpha=0.2, linestyle='--')
    plt.subplots_adjust(bottom=0.25)
    buf = BytesIO(); fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', pad_inches=0.08); plt.close(); buf.seek(0)
    return buf

def chart_donut_sectors():
    fig, ax = plt.subplots(figsize=(3.3, 2.4))
    fig.patch.set_facecolor('white')
    vals = [43.6, 26.5, 15.9, 14.0]
    labels = ['Energy\n(Oil & Gas)', 'Mfg\n(Cement)', 'Mfg\n(Steel)', 'Other']
    colors = ['#1B5E20', '#2E7D32', '#43A047', '#A5D6A7']
    wedges, texts, autos = ax.pie(vals, labels=labels, autopct='%1.1f%%', colors=colors, startangle=90, pctdistance=0.78, textprops={'fontsize':7})
    for a in autos: a.set_fontsize(6); a.set_fontweight('bold')
    ax.add_patch(plt.Circle((0,0), 0.55, fc='white'))
    ax.text(0, 0.06, '22.03M', fontsize=10, fontweight='bold', ha='center', color='#1B3A6B')
    ax.text(0, -0.1, 'tCO\u2082e', fontsize=7, ha='center', color='#666')
    ax.set_title('Financed Emissions by Sector', fontsize=9, fontweight='bold', color='#1B3A6B', pad=4)
    plt.tight_layout(pad=0.3)
    buf = BytesIO(); fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', pad_inches=0.05); plt.close(); buf.seek(0)
    return buf

def chart_scope1_pie():
    fig, ax = plt.subplots(figsize=(3.6, 2.8))
    fig.patch.set_facecolor('white')
    vals = [64.7, 35.3]
    colors = ['#1B3A6B', '#7B1FA2']
    wedges, texts, autos = ax.pie(vals, labels=['Diesel\nCombustion','Natural Gas\nCombustion'], autopct='%1.1f%%', colors=colors, startangle=90, textprops={'fontsize':9}, pctdistance=0.6)
    for a in autos: a.set_color('white'); a.set_fontweight('bold')
    ax.set_title('Scope 1 by Source (3,403 tCO\u2082e)', fontsize=10, fontweight='bold', color='#1B3A6B', pad=6)
    plt.tight_layout(pad=0.4)
    buf = BytesIO(); fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', pad_inches=0.05); plt.close(); buf.seek(0)
    return buf

def chart_sectors_hbar():
    fig, ax = plt.subplots(figsize=(4.6, 2.8))
    fig.patch.set_facecolor('white')
    sectors = ['Other', 'Mfg (Steel)', 'Mfg (Cement)', 'Energy (Oil & Gas)']
    values = [3.1, 3.5, 5.8, 9.6]
    colors = ['#A5D6A7', '#66BB6A', '#43A047', '#1B5E20']
    bars = ax.barh(sectors, values, color=colors, height=0.55)
    for bar, v in zip(bars, values):
        ax.text(v+0.1, bar.get_y()+bar.get_height()/2, f'{v:.1f}M', va='center', fontsize=8, fontweight='bold', color='#1B3A6B')
    ax.set_title('Scope 3 Cat 15 by Sector', fontsize=10, fontweight='bold', color='#1B3A6B', pad=6)
    ax.set_xlabel('M tCO\u2082e', fontsize=8)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.set_xlim(0, 11); ax.grid(axis='x', alpha=0.2, linestyle='--')
    plt.tight_layout(pad=0.4)
    buf = BytesIO(); fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', pad_inches=0.05); plt.close(); buf.seek(0)
    return buf

def chart_intensity():
    fig, ax = plt.subplots(figsize=(4.6, 2.3))
    fig.patch.set_facecolor('white')
    x = np.arange(2)
    totals = [22.98, 22.03]
    intensity = [270.4, 239.4]
    ax.fill_between(x, 0, totals, color='#A5D6A7', alpha=0.5)
    ax.plot(x, totals, 's-', color='#1B5E20', linewidth=2, markersize=6)
    for i, v in enumerate(totals):
        ax.annotate(f'{v:.2f}M', (i,v), xytext=(0,6), textcoords="offset points", ha='center', fontsize=8, fontweight='bold', color='#1B5E20')
    ax.set_ylabel('M tCO\u2082e', fontsize=7, color='#1B5E20')
    ax.set_ylim(0, 28)
    ax2 = ax.twinx()
    ax2.plot(x, intensity, 'o-', color='#FF8F00', linewidth=2.5, markersize=8)
    for i, v in enumerate(intensity):
        ax2.annotate(f'{v:.1f}', (i,v), xytext=(0,8), textcoords="offset points", ha='center', fontsize=9, fontweight='bold', color='#FF8F00')
    ax2.set_ylabel('tCO\u2082e/IDR Bn', fontsize=7, color='#FF8F00')
    ax2.set_ylim(220, 290)
    ax.set_xticks(x); ax.set_xticklabels(['FY2023','FY2024'], fontsize=9)
    ax.set_title('Emissions & Intensity Trend', fontsize=10, fontweight='bold', color='#1B3A6B', pad=6)
    ax.spines['top'].set_visible(False); ax2.spines['top'].set_visible(False)
    ax.grid(axis='y', alpha=0.2, linestyle='--')
    plt.tight_layout(pad=0.4)
    buf = BytesIO(); fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', pad_inches=0.08); plt.close(); buf.seek(0)
    return buf

def chart_social_radar():
    fig, ax = plt.subplots(figsize=(3.6, 3.3), subplot_kw=dict(polar=True))
    fig.patch.set_facecolor('white')
    cats = ['Training\n(49.3 hrs)', 'Female Rep\n(42.32%)', 'New Hires\n(2,595)', 'Female Mgmt\n(26.67%)', 'Turnover\n(8.45%)']
    values = [82, 85, 65, 53, 72]; values += values[:1]
    angles = np.linspace(0, 2*np.pi, len(cats), endpoint=False).tolist(); angles += angles[:1]
    ax.fill(angles, values, color='#7B1FA2', alpha=0.15)
    ax.plot(angles, values, 'o-', color='#7B1FA2', linewidth=2, markersize=4)
    ax.set_xticks(angles[:-1]); ax.set_xticklabels(cats, fontsize=7)
    ax.set_ylim(0, 100); ax.set_yticks([20,40,60,80,100]); ax.set_yticklabels(['20','40','60','80','100'], fontsize=5, color='#999')
    ax.set_title('Social Performance', fontsize=10, fontweight='bold', color='#1B3A6B', pad=12)
    ax.grid(color='#E0E0E0', linewidth=0.5)
    plt.tight_layout(pad=0.6)
    buf = BytesIO(); fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', pad_inches=0.08); plt.close(); buf.seek(0)
    return buf

def chart_roadmap():
    fig, ax = plt.subplots(figsize=(9.6, 4.2))
    fig.patch.set_facecolor('white')
    items = [
        ('Solar+Battery Top 5 Sites (-881 tCO\u2082e)', 0, 2, '#1B5E20'),
        ('Borrower Engagement (Top 20 Emitters)', 1, 3, '#43A047'),
        ('Diversity Target Setting', 0, 1, '#00897B'),
        ('Training Program 50hrs Target', 1, 2, '#26A69A'),
        ('SBTi Commitment Letter', 0, 1, '#FF8F00'),
        ('ESG Committee Formation', 0, 2, '#FFB300'),
        ('Transition Plan Development', 2, 2, '#FFA000'),
        ('PCAF Score \u2192 2.5', 0, 1, '#3D6094'),
        ('Automated Scope 3 Collection', 1, 2, '#5B9BD5'),
    ]
    cats = ['ENV','ENV','SOC','SOC','GOV','GOV','GOV','DATA','DATA']
    for i, (name, start, dur, color) in enumerate(reversed(items)):
        ax.barh(i, dur, left=start, height=0.55, color=color, alpha=0.9, edgecolor='white', linewidth=0.5)
        disp = name if len(name)<38 else name[:35]+"..."
        ax.text(start+dur/2, i, disp, ha='center', va='center', fontsize=6.5, fontweight='bold', color='white')
    ax.set_yticks(range(len(items))); ax.set_yticklabels(list(reversed(cats)), fontsize=7, fontweight='bold')
    quarters = ['Q3 2025','Q4 2025','Q1 2026','Q2 2026','Q3 2026']
    for i, q in enumerate(quarters):
        ax.text(i+0.5, len(items)+0.2, q, ha='center', fontsize=9, fontweight='bold', color='#1B3A6B')
    ax.set_xlim(0, 5); ax.set_ylim(-0.5, len(items)+0.5)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False); ax.spines['left'].set_visible(False)
    ax.grid(axis='x', alpha=0.2, linestyle='--')
    ax.set_xticks(range(6)); ax.set_xticklabels(['']*6)
    ax.set_title('ESG Strategic Roadmap 2025\u20132026', fontsize=12, fontweight='bold', color='#1B3A6B', pad=10)
    plt.tight_layout(pad=0.5)
    buf = BytesIO(); fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', pad_inches=0.08); plt.close(); buf.seek(0)
    return buf


# === TEMPLATE POPULATION LOGIC ===

def replace_text(shape, new_text):
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.text = ""
        if para.runs:
            para.runs[0].text = new_text
        else:
            para.text = new_text
        break

def replace_image(slide, image_name, buf):
    for shape in list(slide.shapes):
        if shape.name == image_name and shape.shape_type == 13:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            sp = shape._element
            sp.getparent().remove(sp)
            slide.shapes.add_picture(buf, left, top, width, height)
            return True
    return False

def main():
    print(f"Loading: {TEMPLATE_PATH}")
    prs = Presentation(TEMPLATE_PATH)
    slides = list(prs.slides)
    print(f"Slides: {len(slides)}")

    # Slide 1
    for shape in slides[0].shapes:
        if shape.name in SLIDE1_TEXT:
            replace_text(shape, SLIDE1_TEXT[shape.name])
    print("Slide 1: charts...")
    replace_image(slides[0], "Image 1", chart_stacked_bar())
    replace_image(slides[0], "Image 2", chart_donut_sectors())

    # Slide 2
    for shape in slides[1].shapes:
        if shape.name in SLIDE2_TEXT:
            replace_text(shape, SLIDE2_TEXT[shape.name])
    print("Slide 2: charts...")
    replace_image(slides[1], "Image 1", chart_scope1_pie())
    replace_image(slides[1], "Image 2", chart_sectors_hbar())
    replace_image(slides[1], "Image 3", chart_intensity())

    # Slide 3
    for shape in slides[2].shapes:
        if shape.name in SLIDE3_TEXT:
            replace_text(shape, SLIDE3_TEXT[shape.name])
    print("Slide 3: chart...")
    replace_image(slides[2], "Image 1", chart_social_radar())

    # Slide 4
    print("Slide 4: roadmap...")
    replace_image(slides[3], "Image 1", chart_roadmap())

    prs.save(OUTPUT_PATH)
    print(f"\nDone! Output: {OUTPUT_PATH}")

if __name__ == "__main__":
    main()
