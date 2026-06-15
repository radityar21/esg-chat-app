"""
ESG Executive Summary — 4-Slide PPTX Demo
Run: pip install python-pptx matplotlib numpy
     python generate_exec_pptx_demo.py

Output: ESG_Executive_Summary_Demo.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np

# Colors
NAVY = RGBColor(0x1B, 0x3A, 0x6B)
BLUE = RGBColor(0x3D, 0x60, 0x94)
TEAL = RGBColor(0x00, 0x89, 0x7B)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
AMBER = RGBColor(0xFF, 0x8F, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY_BG = RGBColor(0xF5, 0xF7, 0xFA)
GREY_TEXT = RGBColor(0x66, 0x66, 0x66)
DARK = RGBColor(0x21, 0x21, 0x21)


def add_header_band(slide, title, subtitle=""):
    """Navy header band at top."""
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.9))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(9), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(0xA8, 0xCC, 0xE0)


def add_kpi_card(slide, x, y, w, h, value, unit, label, trend, color, trend_good=True):
    """Colored KPI card."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()
    # Value
    vb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), w - Inches(0.3), Inches(0.5))
    vp = vb.text_frame.paragraphs[0]
    vp.text = value
    vp.font.size = Pt(24)
    vp.font.bold = True
    vp.font.color.rgb = WHITE
    # Unit
    ub = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.6), w - Inches(0.3), Inches(0.25))
    up = ub.text_frame.paragraphs[0]
    up.text = unit
    up.font.size = Pt(9)
    up.font.color.rgb = RGBColor(0xD0, 0xD8, 0xE8)
    # Label
    lb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.85), w - Inches(0.3), Inches(0.25))
    lp = lb.text_frame.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(8)
    lp.font.bold = True
    lp.font.color.rgb = RGBColor(0xE0, 0xE8, 0xF0)
    # Trend badge
    tb = slide.shapes.add_textbox(x + w - Inches(1.1), y + Inches(0.08), Inches(1.0), Inches(0.25))
    tp = tb.text_frame.paragraphs[0]
    tp.text = trend
    tp.font.size = Pt(9)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(0x69, 0xF0, 0xAE) if trend_good else RGBColor(0xFF, 0x82, 0x82)
    tp.alignment = PP_ALIGN.RIGHT


def make_chart_stacked_bar():
    fig, ax = plt.subplots(figsize=(5.5, 2.8))
    fig.patch.set_facecolor('#F5F7FA')
    ax.set_facecolor('#F5F7FA')
    cats = ['FY2023', 'FY2024']
    s1 = [3324, 3403]
    s2 = [44601, 44614]
    s3 = [22930000, 21976797]
    x = np.arange(2)
    ax.bar(x, s3, 0.5, label='Scope 3 (Financed)', color='#4A7AB5')
    ax.bar(x, s2, 0.5, bottom=s3, label='Scope 2 (Energy)', color='#3D6094')
    ax.bar(x, s1, 0.5, bottom=[a+b for a,b in zip(s3,s2)], label='Scope 1', color='#1B3A6B')
    for i, t in enumerate([sum(x) for x in zip(s1,s2,s3)]):
        ax.annotate(f'{t/1e6:.2f}M', (i, t), textcoords="offset points",
                    xytext=(0,5), ha='center', fontsize=9, fontweight='bold', color='#1B3A6B')
    ax.annotate('▼ 4.07%', (1.3, 20000000), fontsize=10, fontweight='bold', color='#2E7D32',
                bbox=dict(boxstyle='round,pad=0.3', fc='#E8F5E9', ec='#2E7D32'))
    ax.set_xticks(x)
    ax.set_xticklabels(cats, fontsize=10)
    ax.set_title('Total GHG Emissions by Scope (tCO₂e)', fontsize=11, fontweight='bold', color='#1B3A6B')
    ax.legend(fontsize=8)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,p: f'{v/1e6:.0f}M'))
    ax.grid(axis='y', alpha=0.3, linestyle='--')
    plt.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, facecolor='#F5F7FA', bbox_inches='tight')
    plt.close()
    buf.seek(0)
    return buf


def make_chart_donut():
    fig, ax = plt.subplots(figsize=(3.0, 2.8))
    fig.patch.set_facecolor('#F5F7FA')
    ax.set_facecolor('#F5F7FA')
    vals = [99.78, 0.20, 0.02]
    labels = ['Scope 3\nFinanced', 'Scope 2\nEnergy', 'Scope 1\nDirect']
    colors = ['#4A7AB5', '#3D6094', '#1B3A6B']
    wedges, texts, autotexts = ax.pie(vals, labels=labels, autopct='%1.1f%%',
                                       colors=colors, startangle=90, pctdistance=0.78,
                                       textprops={'fontsize': 7})
    circle = plt.Circle((0,0), 0.55, fc='#F5F7FA')
    ax.add_patch(circle)
    ax.text(0, 0.05, '22.03M', fontsize=13, fontweight='bold', ha='center', color='#1B3A6B')
    ax.text(0, -0.15, 'tCO₂e', fontsize=8, ha='center', color='#666')
    ax.set_title('Emission Source', fontsize=10, fontweight='bold', color='#1B3A6B')
    plt.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, facecolor='#F5F7FA', bbox_inches='tight')
    plt.close()
    buf.seek(0)
    return buf


def make_chart_scope1_pie():
    fig, ax = plt.subplots(figsize=(3.0, 2.5))
    fig.patch.set_facecolor('#F5F7FA')
    ax.set_facecolor('#F5F7FA')
    vals = [2580.45, 822.91]
    labels = ['Natural Gas\n75.8%', 'Diesel\n24.2%']
    colors = ['#1B3A6B', '#FF8F00']
    ax.pie(vals, labels=labels, autopct='%1.0f%%', colors=colors, startangle=90, textprops={'fontsize': 8})
    ax.set_title('Scope 1 by Source', fontsize=10, fontweight='bold', color='#1B3A6B')
    plt.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, facecolor='#F5F7FA', bbox_inches='tight')
    plt.close()
    buf.seek(0)
    return buf


def make_chart_sectors():
    fig, ax = plt.subplots(figsize=(4.5, 2.5))
    fig.patch.set_facecolor('#F5F7FA')
    ax.set_facecolor('#F5F7FA')
    sectors = ['Energy Oil & Gas', 'Cement', 'Steel', 'Real Estate', 'Agriculture']
    values = [6.9, 4.2, 2.8, 1.5, 0.9]
    ax.barh(sectors, values, color='#3D6094')
    for i, v in enumerate(values):
        ax.text(v + 0.1, i, f'{v:.1f}M', va='center', fontsize=8, fontweight='bold')
    ax.set_title('Top 5 Sectors — Financed Emissions (M tCO₂e)', fontsize=10, fontweight='bold', color='#1B3A6B')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.set_xlim(0, 8.5)
    plt.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, facecolor='#F5F7FA', bbox_inches='tight')
    plt.close()
    buf.seek(0)
    return buf


def make_chart_intensity():
    fig, ax = plt.subplots(figsize=(4.5, 2.5))
    fig.patch.set_facecolor('#F5F7FA')
    ax.set_facecolor('#F5F7FA')
    years = ['FY2023', 'FY2024']
    totals = [22980000, 22030000]
    intensity = [270.4, 239.4]
    x = np.arange(2)
    ax.bar(x, totals, 0.4, color='#1B3A6B', label='Total (tCO₂e)')
    ax.set_ylabel('tCO₂e', fontsize=8)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,p: f'{v/1e6:.0f}M'))
    ax2 = ax.twinx()
    ax2.plot(x, intensity, 'o-', color='#FF8F00', linewidth=2, markersize=8, label='Intensity')
    ax2.set_ylabel('tCO₂e/IDR Bn', fontsize=8, color='#FF8F00')
    for i, v in enumerate(intensity):
        ax2.annotate(f'{v:.1f}', (i, v), textcoords="offset points", xytext=(0, 10),
                     ha='center', fontsize=9, fontweight='bold', color='#FF8F00')
    ax.set_xticks(x)
    ax.set_xticklabels(years, fontsize=9)
    ax.set_title('Emission Intensity Trend', fontsize=10, fontweight='bold', color='#1B3A6B')
    ax.legend(loc='upper left', fontsize=7)
    ax2.legend(loc='upper right', fontsize=7)
    ax.spines['top'].set_visible(False)
    ax.grid(axis='y', alpha=0.3, linestyle='--')
    plt.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, facecolor='#F5F7FA', bbox_inches='tight')
    plt.close()
    buf.seek(0)
    return buf


def make_chart_social():
    fig, ax = plt.subplots(figsize=(5.0, 2.5))
    fig.patch.set_facecolor('#F5F7FA')
    ax.set_facecolor('#F5F7FA')
    cats = ['Female %', 'Mgmt Female %', 'Training hrs', 'Turnover %']
    fy23 = [41.5, 28.3, 30.2, 8.5]
    fy24 = [42.3, 30.1, 35.2, 7.8]
    x = np.arange(4)
    w = 0.35
    ax.bar(x - w/2, fy23, w, label='FY2023', color='#A8CCE0')
    ax.bar(x + w/2, fy24, w, label='FY2024', color='#1B3A6B')
    ax.set_xticks(x)
    ax.set_xticklabels(cats, fontsize=8)
    ax.set_title('Social Metrics — Year-over-Year', fontsize=10, fontweight='bold', color='#1B3A6B')
    ax.legend(fontsize=8)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(axis='y', alpha=0.3, linestyle='--')
    plt.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, facecolor='#F5F7FA', bbox_inches='tight')
    plt.close()
    buf.seek(0)
    return buf


# =============================================================================
# SLIDE 1: ESG DASHBOARD
# =============================================================================
def slide_1_dashboard(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_band(slide, "ESG Performance Dashboard", "Multi-Framework Sustainability Report — FY2024")

    # KPI cards
    kpis = [
        ("22.03M", "tCO₂e", "Total GHG Emissions", "▼ 4.07%", NAVY, True),
        ("239.4", "tCO₂e/IDR Bn", "Revenue Intensity", "▼ 11.3%", GREEN, True),
        ("3.36", "score (1-5)", "PCAF Data Quality", "▲ 0.12", AMBER, False),
        ("100%", "portfolio", "Coverage", "Stable", RGBColor(0x4A, 0x14, 0x8C), True),
    ]
    for i, (val, unit, label, trend, color, good) in enumerate(kpis):
        x = Inches(0.3) + i * Inches(3.25)
        add_kpi_card(slide, x, Inches(1.1), Inches(3.0), Inches(1.15), val, unit, label, trend, color, good)

    # Charts
    slide.shapes.add_picture(make_chart_stacked_bar(), Inches(0.3), Inches(2.5), Inches(5.8), Inches(3.0))
    slide.shapes.add_picture(make_chart_donut(), Inches(6.3), Inches(2.5), Inches(3.2), Inches(3.0))

    # Scorecard panel (right side)
    panel_x = Inches(9.7)
    sc_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, panel_x, Inches(2.5), Inches(3.3), Inches(1.8))
    sc_bg.fill.solid()
    sc_bg.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    sc_bg.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    tb = slide.shapes.add_textbox(panel_x + Inches(0.15), Inches(2.55), Inches(3.0), Inches(1.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ESG Scorecard"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = NAVY
    scores = [("🟡", "Environmental", "Stable, peers declining faster"),
              ("🟢", "Social", "Female 42.3%, training 35.2 hrs"),
              ("🔴", "Governance", "No SBTi, no ESG committee")]
    for emoji, dim, driver in scores:
        p2 = tf.add_paragraph()
        p2.text = f"{emoji} {dim}"
        p2.font.size = Pt(9)
        p2.font.bold = True
        p2.font.color.rgb = DARK
        p3 = tf.add_paragraph()
        p3.text = f"    {driver}"
        p3.font.size = Pt(7)
        p3.font.color.rgb = GREY_TEXT
        p3.space_after = Pt(2)

    # Board Actions (bottom right)
    ba_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, panel_x, Inches(4.5), Inches(3.3), Inches(2.7))
    ba_bg.fill.solid()
    ba_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE1)
    ba_bg.line.color.rgb = AMBER
    ba_bg.line.width = Pt(1.5)
    tb2 = slide.shapes.add_textbox(panel_x + Inches(0.15), Inches(4.6), Inches(3.0), Inches(2.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "⚡ Board Actions Required"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = AMBER
    actions = [
        "Approve SBTi commitment letter (Q3 2025)",
        "Allocate IDR 15B solar+battery program",
        "Mandate top-20 borrower engagement",
    ]
    for j, act in enumerate(actions):
        pa = tf2.add_paragraph()
        pa.text = f"{j+1}. {act}"
        pa.font.size = Pt(8)
        pa.font.color.rgb = DARK
        pa.space_before = Pt(4)

    # Footer
    ft = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.15), Inches(13.33), Inches(0.35))
    ft.fill.solid()
    ft.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    ft.line.fill.background()
    ftb = slide.shapes.add_textbox(Inches(0.4), Inches(7.18), Inches(12), Inches(0.3))
    ftp = ftb.text_frame.paragraphs[0]
    ftp.text = "AI-Generated | GRI 305 • IFRS S2 • CSRD/ESRS E1 • OJK PSPK | DRAFT | 2026-06-12"
    ftp.font.size = Pt(8)
    ftp.font.color.rgb = GREY_TEXT


# =============================================================================
# SLIDE 2: ENVIRONMENTAL DEEP DIVE
# =============================================================================
def slide_2_environmental(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_band(slide, "Environmental Performance", "GHG Emissions Analysis — Scope 1, 2, 3")

    slide.shapes.add_picture(make_chart_scope1_pie(), Inches(0.3), Inches(1.2), Inches(3.5), Inches(2.8))
    slide.shapes.add_picture(make_chart_sectors(), Inches(4.0), Inches(1.2), Inches(5.0), Inches(2.8))
    slide.shapes.add_picture(make_chart_intensity(), Inches(0.3), Inches(4.2), Inches(5.0), Inches(2.8))

    # Insight box (right side bottom)
    ib = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(4.2),
                                 Inches(4.5), Inches(2.8))
    ib.fill.solid()
    ib.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE1)
    ib.line.color.rgb = AMBER
    tb = slide.shapes.add_textbox(Inches(5.7), Inches(4.35), Inches(4.1), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚠️ KEY INSIGHT"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = AMBER
    p2 = tf.add_paragraph()
    p2.text = "99.78% of emissions from financed portfolio. Operational reduction alone cannot achieve material impact."
    p2.font.size = Pt(9)
    p2.font.color.rgb = DARK
    p2.space_before = Pt(6)
    p3 = tf.add_paragraph()
    p3.text = "\n📊 Diesel dependency (64.7% of Scope 1) is the primary operational lever. Solar+battery at top 5 sites = -881 tCO₂e/yr."
    p3.font.size = Pt(8)
    p3.font.color.rgb = GREY_TEXT
    p3.space_before = Pt(6)
    p4 = tf.add_paragraph()
    p4.text = "\n🏦 Peers: BRI -3.2%, DBS -12% (2023 vs 2022)"
    p4.font.size = Pt(8)
    p4.font.bold = True
    p4.font.color.rgb = RED

    # Quick stat cards (right column top)
    for i, (val, label) in enumerate([("47", "Facilities"), ("2.1/4", "Data Quality"), ("-4.07%", "YoY Change")]):
        y = Inches(1.2) + i * Inches(0.85)
        cb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), y, Inches(1.8), Inches(0.75))
        cb.fill.solid()
        cb.fill.fore_color.rgb = NAVY if i < 2 else GREEN
        cb.line.fill.background()
        ctb = slide.shapes.add_textbox(Inches(9.6), y + Inches(0.05), Inches(1.6), Inches(0.65))
        ctf = ctb.text_frame
        cp = ctf.paragraphs[0]
        cp.text = val
        cp.font.size = Pt(16)
        cp.font.bold = True
        cp.font.color.rgb = WHITE
        cp.alignment = PP_ALIGN.CENTER
        cp2 = ctf.add_paragraph()
        cp2.text = label
        cp2.font.size = Pt(8)
        cp2.font.color.rgb = RGBColor(0xD0, 0xD8, 0xE8)
        cp2.alignment = PP_ALIGN.CENTER


# =============================================================================
# SLIDE 3: SOCIAL & GOVERNANCE
# =============================================================================
def slide_3_social_governance(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_band(slide, "Social & Governance", "Workforce Metrics + ESG Governance Maturity")

    # Social chart (left)
    slide.shapes.add_picture(make_chart_social(), Inches(0.3), Inches(1.2), Inches(5.5), Inches(2.8))

    # E/S/G Column cards (right side)
    col_data = [
        ("E", "Environmental", NAVY, ["Scope 1: 3,403 tCO₂e (+2.4%)", "Scope 2: 44,614 tCO₂e (stable)", "Scope 3: 21.98M tCO₂e (-4.2%)", "Intensity: 239.4 tCO₂e/IDR Bn"]),
        ("S", "Social", TEAL, ["FTE: 12,847 (+3.1%)", "Female: 42.3%", "Training: 35.2 hrs/FTE", "Turnover: 7.8% (down)"]),
        ("G", "Governance", AMBER, ["❌ No ESG Committee", "❌ No SBTi Commitment", "⚠️ No Transition Plan", "✅ PCAF Adopted"]),
    ]
    for i, (letter, title, color, items) in enumerate(col_data):
        x = Inches(6.0) + i * Inches(2.45)
        # Header
        hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.2), Inches(2.3), Inches(0.5))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        htb = slide.shapes.add_textbox(x, Inches(1.22), Inches(2.3), Inches(0.45))
        hp = htb.text_frame.paragraphs[0]
        hp.text = f"{letter} | {title}"
        hp.font.size = Pt(10)
        hp.font.bold = True
        hp.font.color.rgb = WHITE
        hp.alignment = PP_ALIGN.CENTER
        # Items
        itb = slide.shapes.add_textbox(x + Inches(0.1), Inches(1.8), Inches(2.1), Inches(2.5))
        itf = itb.text_frame
        itf.word_wrap = True
        for item in items:
            ip = itf.add_paragraph()
            ip.text = item
            ip.font.size = Pt(8)
            ip.font.color.rgb = DARK
            ip.space_before = Pt(3)

    # Governance gap box (bottom)
    gap_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(4.3),
                                     Inches(12.7), Inches(2.8))
    gap_bg.fill.solid()
    gap_bg.fill.fore_color.rgb = RGBColor(0xFC, 0xE4, 0xEC)
    gap_bg.line.color.rgb = RED
    gtb = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.6))
    gtf = gtb.text_frame
    gtf.word_wrap = True
    gp = gtf.paragraphs[0]
    gp.text = "🔴 Governance Gaps — Framework Compliance"
    gp.font.size = Pt(12)
    gp.font.bold = True
    gp.font.color.rgb = RED
    rows = [
        "IFRS S2 Para 5-9: Board climate oversight — NOT ESTABLISHED",
        "ESRS 2 GOV-1: Sustainability governance body — NOT ESTABLISHED",
        "OJK POJK 51: ESG committee + climate policy — PARTIAL (policy only)",
        "SBTi: Science-based targets — NOT COMMITTED (peers BRI/BCA committed)",
    ]
    for row in rows:
        rp = gtf.add_paragraph()
        rp.text = f"  • {row}"
        rp.font.size = Pt(9)
        rp.font.color.rgb = DARK
        rp.space_before = Pt(3)


# =============================================================================
# SLIDE 4: STRATEGIC ROADMAP
# =============================================================================
def slide_4_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_band(slide, "Strategic Roadmap", "Priority Actions & Implementation Timeline")

    # Timeline header (quarters)
    quarters = ["Q3 2025", "Q4 2025", "Q1 2026", "Q2 2026", "Q3 2026"]
    q_width = Inches(2.2)
    for i, q in enumerate(quarters):
        x = Inches(2.5) + i * q_width
        qtb = slide.shapes.add_textbox(x, Inches(1.1), q_width, Inches(0.3))
        qp = qtb.text_frame.paragraphs[0]
        qp.text = q
        qp.font.size = Pt(10)
        qp.font.bold = True
        qp.font.color.rgb = NAVY
        qp.alignment = PP_ALIGN.CENTER
        # Vertical line
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.4), Pt(1), Inches(5.5))
        ln.fill.solid()
        ln.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        ln.line.fill.background()

    # Category labels (left)
    categories = [
        ("ENVIRON.", NAVY, Inches(1.6)),
        ("SOCIAL", TEAL, Inches(3.2)),
        ("GOVERN.", AMBER, Inches(4.8)),
        ("DATA", BLUE, Inches(6.0)),
    ]
    for label, color, y in categories:
        lbl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.2), y, Inches(1.8), Inches(0.4))
        lbl.fill.solid()
        lbl.fill.fore_color.rgb = color
        lbl.line.fill.background()
        ltb = slide.shapes.add_textbox(Inches(0.25), y + Inches(0.02), Inches(1.7), Inches(0.35))
        lp = ltb.text_frame.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(9)
        lp.font.bold = True
        lp.font.color.rgb = WHITE
        lp.alignment = PP_ALIGN.CENTER

    # Timeline bars (Gantt-style)
    bars = [
        # (category_y, start_q, duration_q, text, color)
        (Inches(1.6), 0, 2, "Solar+Battery Top 5 Sites (-881 tCO₂e)", NAVY),
        (Inches(2.1), 1, 3, "Borrower Engagement Program (Top 20 Emitters)", BLUE),
        (Inches(3.2), 0, 1, "Diversity Target Setting", TEAL),
        (Inches(3.6), 1, 2, "Training Program Expansion (40hrs target)", TEAL),
        (Inches(4.8), 0, 1, "SBTi Commitment Letter", AMBER),
        (Inches(5.2), 0, 2, "ESG Committee Formation", AMBER),
        (Inches(5.6), 2, 2, "Transition Plan Development", AMBER),
        (Inches(6.0), 0, 1, "PCAF Score → 2.5 (Data Improvement)", BLUE),
        (Inches(6.4), 1, 2, "Automated Scope 3 Data Collection", BLUE),
    ]
    for y, start, dur, text, color in bars:
        x = Inches(2.5) + start * q_width + Inches(0.1)
        w = dur * q_width - Inches(0.2)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.35))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        btb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.02), w - Inches(0.2), Inches(0.3))
        bp = btb.text_frame.paragraphs[0]
        bp.text = text
        bp.font.size = Pt(7)
        bp.font.bold = True
        bp.font.color.rgb = WHITE


# =============================================================================
# MAIN
# =============================================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_1_dashboard(prs)
    slide_2_environmental(prs)
    slide_3_social_governance(prs)
    slide_4_roadmap(prs)

    out = "ESG_Executive_Summary_Demo.pptx"
    prs.save(out)
    print(f"Done! Generated: {out}")


if __name__ == "__main__":
    main()
